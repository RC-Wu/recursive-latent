#!/usr/bin/env python3
"""Compose editable PPTX method figures for the PS-RSLE paper revision."""

from __future__ import annotations

import json
from pathlib import Path

from PIL import Image, ImageChops
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_LINE_DASH_STYLE
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE, MSO_CONNECTOR
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path("/Users/fanta/code/agent/Code/recursive_3d_generative_growth")
OUT_DIR = ROOT / "paper_siga/figures/method_diagram_pptx_20260512"

INK = "1B232B"
MUTED = "5A6572"
LINE = "CFD7E3"
LANE = "F5F8FB"
BLUE = "2F6FDB"
BLUE_L = "EAF2FF"
TEAL = "168A83"
TEAL_L = "E7F5F3"
GREEN = "2E7D4F"
GREEN_L = "EAF6ED"
AMBER = "B76E00"
AMBER_L = "FFF4E2"
ROSE = "B84A55"
ROSE_L = "FBECEF"
VIOLET = "6A5EA8"
VIOLET_L = "F0EEFA"
GRAY_L = "F2F4F7"
WHITE = "FFFFFF"
SUBTLE_BLUE = "F8FAFD"
SUBTLE_ROSE = "FDF4F5"

ASSET_DIR = OUT_DIR / "component_assets"
LOCAL_ASSETS = {
    "tree_input": ROOT / "visuals/trellis2_mesh_first_reconstruct_20260507_1745/ifs_branch_tree_r512/input_normalized_preview.png",
    "tree_slat": ROOT / "visuals/trellis2_mesh_first_reconstruct_20260507_1745/ifs_branch_tree_r512/shape_slat_reconstruct_preview.png",
    "tree_roundtrip": ROOT / "visuals/trellis2_mesh_first_reconstruct_20260507_1745/ifs_branch_tree_r512/ovoxel_roundtrip_preview.png",
    "vine_overview": ROOT / "visuals/matched_camera_zoom_existing_roots_20260510/vine_stage5/overview_raw.png",
    "vine_zoom_01": ROOT / "visuals/matched_camera_zoom_existing_roots_20260510/vine_stage5/zoom_01.png",
    "vine_zoom_02": ROOT / "visuals/matched_camera_zoom_existing_roots_20260510/vine_stage5/zoom_02.png",
}


def rgb(hex_color: str) -> RGBColor:
    hex_color = hex_color.strip("#")
    return RGBColor(int(hex_color[0:2], 16), int(hex_color[2:4], 16), int(hex_color[4:6], 16))


def inch(value: float) -> int:
    return Inches(value)


def prepare_component_assets() -> dict[str, Path]:
    """Create cropped transparent component images used inside editable PPTX diagrams."""
    ASSET_DIR.mkdir(parents=True, exist_ok=True)

    def transparent_crop(src: Path, name: str, *, drop_top: int = 0, threshold: int = 246, pad: int = 20) -> Path:
        out = ASSET_DIR / f"{name}.png"
        im = Image.open(src).convert("RGBA")
        if drop_top:
            im = im.crop((0, drop_top, im.width, im.height))

        # Convert near-white rendered backgrounds to transparent while keeping colored geometry.
        data = im.getdata()
        converted = []
        for r, g, b, a in data:
            if a == 0:
                converted.append((r, g, b, 0))
            elif r >= threshold and g >= threshold and b >= threshold:
                converted.append((r, g, b, 0))
            else:
                converted.append((r, g, b, a))
        im.putdata(converted)

        alpha = im.getchannel("A")
        bbox = alpha.getbbox()
        if bbox:
            left, upper, right, lower = bbox
            left = max(0, left - pad)
            upper = max(0, upper - pad)
            right = min(im.width, right + pad)
            lower = min(im.height, lower + pad)
            im = im.crop((left, upper, right, lower))
        im.save(out)
        return out

    def alpha_crop(src: Path, name: str, *, pad: int = 14) -> Path:
        out = ASSET_DIR / f"{name}.png"
        im = Image.open(src).convert("RGBA")
        alpha = im.getchannel("A")
        bbox = alpha.getbbox()
        if bbox:
            left, upper, right, lower = bbox
            left = max(0, left - pad)
            upper = max(0, upper - pad)
            right = min(im.width, right + pad)
            lower = min(im.height, lower + pad)
            im = im.crop((left, upper, right, lower))
        im.save(out)
        return out

    return {
        "tree_input": transparent_crop(LOCAL_ASSETS["tree_input"], "tree_input_transparent", drop_top=145),
        "tree_slat": transparent_crop(LOCAL_ASSETS["tree_slat"], "tree_slat_transparent", drop_top=145),
        "tree_roundtrip": transparent_crop(LOCAL_ASSETS["tree_roundtrip"], "tree_roundtrip_transparent", drop_top=145),
        "vine_overview": alpha_crop(LOCAL_ASSETS["vine_overview"], "vine_overview_crop"),
        "vine_zoom_01": alpha_crop(LOCAL_ASSETS["vine_zoom_01"], "vine_zoom_01_crop"),
        "vine_zoom_02": alpha_crop(LOCAL_ASSETS["vine_zoom_02"], "vine_zoom_02_crop"),
    }


def add_image_crop(slide, path: Path, x: float, y: float, w: float, h: float, *, crop_left=0.0, crop_right=0.0, crop_top=0.0, crop_bottom=0.0):
    pic = slide.shapes.add_picture(str(path), inch(x), inch(y), inch(w), inch(h))
    pic.crop_left = crop_left
    pic.crop_right = crop_right
    pic.crop_top = crop_top
    pic.crop_bottom = crop_bottom
    return pic


def set_slide_size(prs: Presentation, width: float, height: float) -> None:
    prs.slide_width = inch(width)
    prs.slide_height = inch(height)


def blank_slide(prs: Presentation):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = rgb(WHITE)
    return slide




def disable_shadow(shape) -> None:
    try:
        shape.shadow.inherit = False
    except Exception:
        pass

def add_text(
    slide,
    text: str,
    x: float,
    y: float,
    w: float,
    h: float,
    *,
    size: float = 8.0,
    color: str = INK,
    bold: bool = False,
    align=PP_ALIGN.CENTER,
    font: str = "Arial",
    valign=MSO_ANCHOR.MIDDLE,
) -> None:
    shape = slide.shapes.add_textbox(inch(x), inch(y), inch(w), inch(h))
    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = inch(0.03)
    tf.margin_right = inch(0.03)
    tf.margin_top = inch(0.01)
    tf.margin_bottom = inch(0.01)
    tf.vertical_anchor = valign
    lines = text.split("\n")
    for idx, line in enumerate(lines):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.alignment = align
        p.space_after = Pt(0)
        p.space_before = Pt(0)
        run = p.add_run()
        run.text = line
        run.font.name = font
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = rgb(color)


def add_box(
    slide,
    x: float,
    y: float,
    w: float,
    h: float,
    title: str,
    body: str = "",
    *,
    fill: str = WHITE,
    line: str = LINE,
    title_color: str = INK,
    body_color: str = MUTED,
    title_size: float = 8.0,
    body_size: float = 6.2,
    bold: bool = True,
    radius_shape=MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
    dash: bool = False,
):
    shape = slide.shapes.add_shape(radius_shape, inch(x), inch(y), inch(w), inch(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb(fill)
    shape.line.color.rgb = rgb(line)
    shape.line.width = Pt(0.8)
    disable_shadow(shape)
    if dash:
        shape.line.dash_style = MSO_LINE_DASH_STYLE.DASH
    tf = shape.text_frame
    tf.clear()
    tf.margin_left = inch(0.07)
    tf.margin_right = inch(0.07)
    tf.margin_top = inch(0.04)
    tf.margin_bottom = inch(0.04)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.space_after = Pt(1)
    r = p.add_run()
    r.text = title
    r.font.name = "Arial"
    r.font.size = Pt(title_size)
    r.font.bold = bold
    r.font.color.rgb = rgb(title_color)
    if body:
        for line_text in body.split("\n"):
            p = tf.add_paragraph()
            p.alignment = PP_ALIGN.CENTER
            p.space_before = Pt(0)
            p.space_after = Pt(0)
            r = p.add_run()
            r.text = line_text
            r.font.name = "Arial"
            r.font.size = Pt(body_size)
            r.font.bold = False
            r.font.color.rgb = rgb(body_color)
    return shape


def add_panel(
    slide,
    x: float,
    y: float,
    w: float,
    h: float,
    *,
    fill: str,
    line: str,
    width: float = 0.75,
    radius_shape=MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
):
    shape = slide.shapes.add_shape(radius_shape, inch(x), inch(y), inch(w), inch(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb(fill)
    shape.line.color.rgb = rgb(line)
    shape.line.width = Pt(width)
    disable_shadow(shape)
    return shape


def add_arrow(slide, x1: float, y1: float, x2: float, y2: float, *, color: str = MUTED, width: float = 1.0, dash: bool = False):
    connector = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, inch(x1), inch(y1), inch(x2), inch(y2))
    connector.line.color.rgb = rgb(color)
    connector.line.width = Pt(width)
    connector.line.end_arrowhead = True
    if dash:
        connector.line.dash_style = MSO_LINE_DASH_STYLE.DASH
    return connector


def add_elbow_arrow(
    slide,
    points: list[tuple[float, float]],
    *,
    color: str = MUTED,
    width: float = 0.9,
    dash: bool = False,
) -> None:
    for idx, (a, b) in enumerate(zip(points, points[1:])):
        line = add_arrow(slide, a[0], a[1], b[0], b[1], color=color, width=width, dash=dash) if idx == len(points) - 2 else add_line(slide, a[0], a[1], b[0], b[1], color=color, width=width, dash=dash)
        if idx < len(points) - 2:
            line.line.end_arrowhead = False


def add_line(slide, x1: float, y1: float, x2: float, y2: float, *, color: str = LINE, width: float = 0.7, dash: bool = False):
    line = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, inch(x1), inch(y1), inch(x2), inch(y2))
    line.line.color.rgb = rgb(color)
    line.line.width = Pt(width)
    if dash:
        line.line.dash_style = MSO_LINE_DASH_STYLE.DASH
    return line


def add_image_contain(slide, path: Path, x: float, y: float, w: float, h: float, *, transparency: int | None = None):
    with Image.open(path) as im:
        iw, ih = im.size
    scale = min(w / iw, h / ih)
    ww = iw * scale
    hh = ih * scale
    opts = {"path": str(path), "x": inch(x + (w - ww) / 2), "y": inch(y + (h - hh) / 2), "w": inch(ww), "h": inch(hh)}
    if transparency is not None:
        opts["transparency"] = transparency
    return slide.shapes.add_picture(opts["path"], opts["x"], opts["y"], opts["w"], opts["h"])


def add_stage_number(slide, n: str, x: float, y: float, color: str) -> None:
    circ = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, inch(x), inch(y), inch(0.24), inch(0.24))
    circ.fill.solid()
    circ.fill.fore_color.rgb = rgb(color)
    circ.line.color.rgb = rgb(WHITE)
    circ.line.width = Pt(0.9)
    disable_shadow(circ)
    add_text(slide, n, x + 0.02, y + 0.015, 0.20, 0.17, size=6.2, color=WHITE, bold=True)


def add_role_pill(slide, text: str, x: float, y: float, color: str, w: float = 0.74) -> None:
    pill = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, inch(x), inch(y), inch(w), inch(0.18))
    pill.fill.solid()
    pill.fill.fore_color.rgb = rgb(WHITE)
    pill.line.color.rgb = rgb(color)
    pill.line.width = Pt(0.65)
    disable_shadow(pill)
    add_text(slide, text, x + 0.02, y + 0.025, w - 0.04, 0.11, size=4.8, color=color, bold=True)


def add_sparse_points(slide, x: float, y: float, *, scale: float = 1.0, accent: str = BLUE, muted: str = TEAL) -> None:
    coords = [
        (0.04, 0.19, accent), (0.18, 0.08, muted), (0.31, 0.18, accent), (0.47, 0.08, muted),
        (0.13, 0.34, accent), (0.29, 0.33, muted), (0.43, 0.34, accent), (0.57, 0.27, muted),
        (0.23, 0.52, muted), (0.39, 0.50, accent), (0.54, 0.49, muted), (0.70, 0.42, accent),
        (0.07, 0.62, muted), (0.27, 0.72, accent), (0.50, 0.69, accent), (0.72, 0.64, muted),
    ]
    for cx, cy, col in coords:
        shape = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
            inch(x + cx * scale),
            inch(y + cy * scale),
            inch(0.055 * scale),
            inch(0.055 * scale),
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = rgb(col)
        shape.line.color.rgb = rgb(WHITE)
        shape.line.width = Pt(0.2)
    for a, b in [((0.18, 0.08), (0.31, 0.18)), ((0.31, 0.18), (0.43, 0.34)), ((0.29, 0.33), (0.39, 0.50)), ((0.50, 0.69), (0.72, 0.64)), ((0.13, 0.34), (0.27, 0.72))]:
        add_line(slide, x + a[0] * scale + 0.03 * scale, y + a[1] * scale + 0.03 * scale,
                 x + b[0] * scale + 0.03 * scale, y + b[1] * scale + 0.03 * scale,
                 color="B8C4D2", width=0.45)


def add_condition_icons(slide, x: float, y: float) -> None:
    img = add_box(slide, x, y, 0.23, 0.18, "", fill=BLUE_L, line=BLUE, radius_shape=MSO_AUTO_SHAPE_TYPE.RECTANGLE)
    tri = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ISOSCELES_TRIANGLE, inch(x + 0.03), inch(y + 0.08), inch(0.07), inch(0.07))
    tri.fill.solid()
    tri.fill.fore_color.rgb = rgb(TEAL)
    tri.line.color.rgb = rgb(TEAL)
    disable_shadow(tri)
    sun = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, inch(x + 0.15), inch(y + 0.04), inch(0.035), inch(0.035))
    sun.fill.solid()
    sun.fill.fore_color.rgb = rgb(AMBER)
    sun.line.color.rgb = rgb(AMBER)
    disable_shadow(sun)
    add_text(slide, "T", x + 0.34, y + 0.00, 0.18, 0.18, size=7.5, bold=True, color=TEAL)
    merge = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.FLOWCHART_SUMMING_JUNCTION, inch(x + 0.62), inch(y + 0.04), inch(0.11), inch(0.11))
    merge.fill.solid()
    merge.fill.fore_color.rgb = rgb(GREEN_L)
    merge.line.color.rgb = rgb(GREEN)
    merge.line.width = Pt(0.75)
    disable_shadow(merge)
    empty = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, inch(x + 0.98), inch(y + 0.045), inch(0.09), inch(0.09))
    empty.fill.background()
    empty.line.color.rgb = rgb(MUTED)
    empty.line.width = Pt(0.75)
    disable_shadow(empty)


def add_asset_icon(slide, x: float, y: float, *, color: str = AMBER) -> None:
    base = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.HEXAGON, inch(x), inch(y), inch(0.38), inch(0.30))
    base.fill.solid()
    base.fill.fore_color.rgb = rgb(AMBER_L if color == AMBER else BLUE_L)
    base.line.color.rgb = rgb(color)
    base.line.width = Pt(0.8)
    top = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ISOSCELES_TRIANGLE, inch(x + 0.09), inch(y - 0.06), inch(0.20), inch(0.18))
    top.fill.solid()
    top.fill.fore_color.rgb = rgb(color)
    top.line.color.rgb = rgb(color)


def build_sparse_interface(assets: dict[str, Path]) -> Path:
    prs = Presentation()
    set_slide_size(prs, 7.30, 2.70)
    slide = blank_slide(prs)

    add_text(slide, "Sparse Latent Generator Interface", 0.18, 0.06, 3.30, 0.18, size=8.4, bold=True, align=PP_ALIGN.LEFT)
    add_line(slide, 0.18, 0.30, 7.12, 0.30, color=LINE, width=0.50)

    # Condition as a quiet broadcast lane. Keep labels large enough for one-column print.
    cond = add_panel(slide, 0.28, 0.45, 1.48, 0.34, fill="F8FBFF", line=BLUE, width=0.55, radius_shape=MSO_AUTO_SHAPE_TYPE.RECTANGLE)
    add_text(slide, "Condition y", 0.39, 0.50, 0.72, 0.11, size=6.4, bold=True, color=BLUE, align=PP_ALIGN.LEFT)
    add_condition_icons(slide, 1.03, 0.52)
    add_text(slide, "image / text / category / empty", 0.39, 0.64, 1.02, 0.10, size=4.3, color=MUTED, align=PP_ALIGN.LEFT)
    add_line(slide, 1.76, 0.62, 6.83, 0.62, color=BLUE, width=0.60)

    # Optional input asset.
    in_panel = add_panel(slide, 0.28, 1.28, 0.74, 0.52, fill=WHITE, line=LINE, width=0.50)
    add_image_crop(slide, assets["tree_input"], 0.35, 1.33, 0.60, 0.31, crop_right=0.24)
    add_text(slide, "optional input asset\nx_in", 0.20, 1.84, 0.94, 0.22, size=5.3, color=MUTED)

    enc = add_box(slide, 1.28, 1.26, 0.78, 0.42, "Encoder", "Enc_θ", fill=BLUE_L, line=BLUE, title_size=6.7, body_size=5.7, title_color=BLUE, radius_shape=MSO_AUTO_SHAPE_TYPE.RECTANGLE)

    # Central sparse latent state.
    state = add_panel(slide, 2.26, 0.98, 1.66, 1.00, fill=TEAL_L, line=TEAL, width=0.85, radius_shape=MSO_AUTO_SHAPE_TYPE.RECTANGLE)
    add_text(slide, "Sparse latent state", 2.43, 1.07, 1.32, 0.12, size=6.2, bold=True, color=TEAL)
    add_text(slide, "u=(V,F)", 2.54, 1.23, 1.10, 0.18, size=8.8, bold=True, color=INK)
    add_image_crop(slide, assets["tree_slat"], 2.43, 1.43, 0.82, 0.36, crop_right=0.30)
    add_sparse_points(slide, 3.16, 1.31, scale=0.44, accent=TEAL, muted=BLUE)
    add_box(slide, 2.20, 2.08, 0.76, 0.22, "V: sparse support", fill=WHITE, line=TEAL, title_size=5.0, title_color=MUTED, radius_shape=MSO_AUTO_SHAPE_TYPE.RECTANGLE)
    add_box(slide, 3.20, 2.08, 0.76, 0.22, "F: latent features", fill=WHITE, line=TEAL, title_size=5.0, title_color=MUTED, radius_shape=MSO_AUTO_SHAPE_TYPE.RECTANGLE)
    add_line(slide, 2.64, 1.98, 2.64, 2.07, color=TEAL, width=0.50)
    add_line(slide, 3.57, 1.98, 3.57, 2.07, color=TEAL, width=0.50)

    sampler = add_box(
        slide,
        4.30,
        1.15,
        1.34,
        0.58,
        "Native sparse-latent\nsampler",
        "ū = 𝒩_θ(ũ; y, ε)",
        fill=GREEN_L,
        line=GREEN,
        title_size=5.9,
        body_size=6.0,
        title_color=GREEN,
        radius_shape=MSO_AUTO_SHAPE_TYPE.RECTANGLE,
    )
    for px in [4.46, 4.67, 4.88, 5.09]:
        circ = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, inch(px), inch(1.56), inch(0.050), inch(0.050))
        circ.fill.solid()
        circ.fill.fore_color.rgb = rgb(GREEN)
        circ.line.color.rgb = rgb(WHITE)
        disable_shadow(circ)
    add_arrow(slide, 4.44, 1.585, 5.21, 1.585, color=GREEN, width=0.50, dash=True)

    latent = add_box(slide, 5.86, 1.22, 0.58, 0.46, "ū", "sampled /\nmodel-consistent", fill=WHITE, line=GREEN, title_size=8.4, body_size=4.4, title_color=GREEN, radius_shape=MSO_AUTO_SHAPE_TYPE.RECTANGLE)
    dec = add_box(slide, 6.58, 1.24, 0.58, 0.42, "Decoder", "Dec_θ", fill=AMBER_L, line=AMBER, title_size=5.9, body_size=5.1, title_color=AMBER, radius_shape=MSO_AUTO_SHAPE_TYPE.RECTANGLE)
    out_panel = add_panel(slide, 6.46, 1.90, 0.70, 0.44, fill=WHITE, line=LINE, width=0.50)
    add_image_crop(slide, assets["tree_roundtrip"], 6.50, 1.94, 0.62, 0.30, crop_right=0.24)
    add_text(slide, "3D asset\nx_out", 6.32, 2.34, 0.98, 0.22, size=5.2, color=MUTED)

    add_arrow(slide, 1.02, 1.52, 1.26, 1.52, color=MUTED, width=0.78)
    add_arrow(slide, 2.06, 1.50, 2.24, 1.50, color=MUTED, width=0.78)
    add_arrow(slide, 3.92, 1.50, 4.28, 1.50, color=MUTED, width=0.78)
    add_arrow(slide, 5.64, 1.50, 5.84, 1.50, color=MUTED, width=0.78)
    add_arrow(slide, 6.44, 1.50, 6.56, 1.50, color=MUTED, width=0.78)
    add_arrow(slide, 6.86, 1.66, 6.82, 1.90, color=AMBER, width=0.65)

    for x, y2 in [(1.67, 1.25), (4.98, 1.13), (6.86, 1.22)]:
        add_arrow(slide, x, 0.62, x, y2, color=BLUE, width=0.58)

    note = add_box(slide, 4.33, 1.94, 1.52, 0.24, "later specialization: 𝒩_θ(ũ; m,y,ε)", fill=WHITE, line=LINE, title_size=4.9, title_color=MUTED, dash=True, radius_shape=MSO_AUTO_SHAPE_TYPE.RECTANGLE)
    note.text_frame.margin_left = inch(0.03)

    out = OUT_DIR / "sparse_latent_generator_interface_20260512.pptx"
    prs.save(out)
    return out


def add_depth_chip(slide, text: str, x: float, y: float, color: str) -> None:
    add_box(slide, x, y, 0.58, 0.20, text, fill=WHITE, line=color, title_size=5.6, title_color=color)


def add_handle_miniature(slide, x: float, y: float) -> None:
    add_sparse_points(slide, x, y, scale=0.52, accent=BLUE, muted=TEAL)
    for cx, cy, col in [(0.07, 0.02, ROSE), (0.30, 0.18, AMBER), (0.22, 0.36, GREEN)]:
        circ = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, inch(x + cx), inch(y + cy), inch(0.08), inch(0.08))
        circ.fill.solid()
        circ.fill.fore_color.rgb = rgb(col)
        circ.line.color.rgb = rgb(WHITE)
        circ.line.width = Pt(0.4)


def build_method_overview(assets: dict[str, Path]) -> Path:
    prs = Presentation()
    set_slide_size(prs, 13.33, 6.90)
    slide = blank_slide(prs)

    add_text(slide, "Projection-Stabilized Recursive Sparse-Latent Execution", 0.30, 0.12, 6.55, 0.24, size=12.4, bold=True, align=PP_ALIGN.LEFT)
    add_text(slide, "one recursive depth step", 10.30, 0.17, 2.70, 0.18, size=7.1, color=MUTED, align=PP_ALIGN.RIGHT)
    add_line(slide, 0.30, 0.46, 13.03, 0.46, color=LINE, width=0.55)

    # Fixed generator substrate is a slim reference lane; the recursive transition carries the main story.
    outer = add_panel(slide, 0.42, 0.72, 12.49, 0.66, fill="FAFBFD", line="E2E8F0", width=0.45, radius_shape=MSO_AUTO_SHAPE_TYPE.RECTANGLE)
    add_text(slide, "fixed sparse-latent generator substrate", 0.62, 0.82, 2.52, 0.13, size=6.4, bold=True, color=BLUE, align=PP_ALIGN.LEFT)
    add_box(slide, 0.92, 1.04, 0.64, 0.23, "Enc_θ", "", fill=WHITE, line=BLUE, title_size=6.4, title_color=BLUE, radius_shape=MSO_AUTO_SHAPE_TYPE.RECTANGLE)
    add_box(slide, 2.22, 1.00, 0.98, 0.31, "u=(V,F)", "", fill=WHITE, line=TEAL, title_size=6.6, title_color=TEAL, radius_shape=MSO_AUTO_SHAPE_TYPE.RECTANGLE)
    add_box(slide, 4.20, 1.00, 1.42, 0.31, "𝒩_θ", "sampler", fill=WHITE, line=GREEN, title_size=6.8, body_size=4.2, title_color=GREEN, radius_shape=MSO_AUTO_SHAPE_TYPE.RECTANGLE)
    add_box(slide, 6.74, 1.04, 0.64, 0.23, "Dec_θ", "", fill=WHITE, line=AMBER, title_size=6.4, title_color=AMBER, radius_shape=MSO_AUTO_SHAPE_TYPE.RECTANGLE)
    asset_card = add_panel(slide, 8.92, 0.93, 1.06, 0.34, fill=WHITE, line=LINE, width=0.45, radius_shape=MSO_AUTO_SHAPE_TYPE.RECTANGLE)
    add_image_crop(slide, assets["tree_roundtrip"], 9.16, 0.97, 0.54, 0.22, crop_right=0.22)
    for a, b in [(1.56, 2.20), (3.20, 4.18), (5.62, 6.72), (7.38, 8.90)]:
        add_arrow(slide, a, 1.16, b, 1.16, color="7E8A98", width=0.62)
    add_text(slide, "codec + native sparse-latent update reused below", 10.10, 1.03, 2.58, 0.15, size=6.3, color=MUTED, align=PP_ALIGN.RIGHT)

    add_text(slide, "PS-RSLE transition at depth d", 0.48, 1.66, 2.80, 0.18, size=7.8, bold=True, color=TEAL, align=PP_ALIGN.LEFT)
    add_line(slide, 0.48, 1.93, 12.88, 1.93, color="E4E9F0", width=0.55)

    # Stage 1: current executable state.
    state = add_panel(slide, 0.56, 2.18, 2.22, 2.38, fill=BLUE_L, line=BLUE, width=0.85, radius_shape=MSO_AUTO_SHAPE_TYPE.RECTANGLE)
    add_stage_number(slide, "1", 0.70, 2.32, BLUE)
    add_text(slide, "Executable state", 1.03, 2.35, 1.32, 0.15, size=8.4, bold=True, color=BLUE, align=PP_ALIGN.LEFT)
    add_text(slide, "s_d=(u_d,A_d)", 0.98, 2.66, 1.36, 0.17, size=7.8, bold=True, color=INK)
    add_image_crop(slide, assets["tree_slat"], 0.90, 2.98, 1.08, 0.52, crop_right=0.30)
    add_handle_miniature(slide, 1.62, 3.17)
    add_text(slide, "rule-readable active handles", 0.86, 3.88, 1.60, 0.14, size=6.1, color=MUTED)
    add_text(slide, "later rules can read only active, root-reachable support", 0.80, 4.10, 1.72, 0.26, size=5.6, color=MUTED)

    # Stage 2+3: proposal and admission as one compact logic cluster with two stacked boxes.
    prop = add_panel(slide, 3.20, 2.14, 2.34, 1.12, fill=TEAL_L, line=TEAL, width=0.85, radius_shape=MSO_AUTO_SHAPE_TYPE.RECTANGLE)
    add_stage_number(slide, "2", 3.34, 2.29, TEAL)
    add_text(slide, "Rule proposal", 3.70, 2.32, 1.30, 0.15, size=8.4, bold=True, color=TEAL, align=PP_ALIGN.LEFT)
    add_text(slide, "ρ(h_i;s_d) → P_i", 3.64, 2.62, 1.34, 0.16, size=7.2, bold=True, color=INK, align=PP_ALIGN.LEFT)
    add_text(slide, "tentative sparse edit", 3.64, 2.84, 1.30, 0.13, size=6.2, color=MUTED, align=PP_ALIGN.LEFT)
    add_sparse_points(slide, 4.92, 2.55, scale=0.34, accent=TEAL, muted=AMBER)

    ctrl = add_panel(slide, 3.20, 3.58, 2.34, 1.10, fill=AMBER_L, line=AMBER, width=0.85, radius_shape=MSO_AUTO_SHAPE_TYPE.RECTANGLE)
    add_stage_number(slide, "3", 3.34, 3.72, AMBER)
    add_text(slide, "Admission + control", 3.70, 3.75, 1.50, 0.15, size=8.1, bold=True, color=AMBER, align=PP_ALIGN.LEFT)
    add_text(slide, "P_d → C_d", 3.72, 4.05, 0.92, 0.16, size=7.4, bold=True, color=INK, align=PP_ALIGN.LEFT)
    add_text(slide, "attach · mask · budget", 3.72, 4.28, 1.38, 0.13, size=6.2, color=MUTED, align=PP_ALIGN.LEFT)
    add_arrow(slide, 4.37, 3.26, 4.37, 3.56, color=AMBER, width=0.90)

    # Stage 4: sample/decode.
    sample = add_panel(slide, 6.18, 2.60, 2.42, 1.72, fill=GREEN_L, line=GREEN, width=0.85, radius_shape=MSO_AUTO_SHAPE_TYPE.RECTANGLE)
    add_stage_number(slide, "4", 6.32, 2.75, GREEN)
    add_text(slide, "Sample + decode", 6.68, 2.78, 1.34, 0.15, size=8.4, bold=True, color=GREEN, align=PP_ALIGN.LEFT)
    add_text(slide, "ū = 𝒩_θ(ũ; C_d,y,ε_d)", 6.56, 3.08, 1.68, 0.15, size=7.0, color=INK, align=PP_ALIGN.LEFT)
    add_text(slide, "x' = Dec_θ(ū,y)", 6.56, 3.30, 1.25, 0.15, size=7.0, color=INK, align=PP_ALIGN.LEFT)
    add_image_crop(slide, assets["vine_zoom_01"], 6.76, 3.55, 1.03, 0.56, crop_left=0.08, crop_right=0.08)
    add_text(slide, "candidate, not committed", 6.58, 4.12, 1.58, 0.13, size=6.1, color=MUTED)

    # Stage 5: projection-centered event.
    proj = add_panel(slide, 9.05, 2.12, 3.78, 2.84, fill=SUBTLE_ROSE, line=ROSE, width=0.90, radius_shape=MSO_AUTO_SHAPE_TYPE.RECTANGLE)
    add_stage_number(slide, "5", 9.21, 2.27, ROSE)
    add_text(slide, "Projection + commit", 9.57, 2.30, 1.62, 0.16, size=8.8, bold=True, color=ROSE, align=PP_ALIGN.LEFT)
    add_text(slide, "keep root-attached active support", 10.94, 2.31, 1.60, 0.13, size=6.0, color=ROSE, align=PP_ALIGN.RIGHT)

    before = add_panel(slide, 9.22, 2.90, 1.05, 0.94, fill=WHITE, line="E0B8BE", width=0.55, radius_shape=MSO_AUTO_SHAPE_TYPE.RECTANGLE)
    add_image_crop(slide, assets["vine_zoom_01"], 9.30, 2.98, 0.86, 0.58, crop_left=0.08, crop_right=0.08)
    for px, py in [(9.98, 3.05), (10.06, 3.18), (9.95, 3.35)]:
        bad = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, inch(px), inch(py), inch(0.060), inch(0.060))
        bad.fill.solid()
        bad.fill.fore_color.rgb = rgb("A8AFB8")
        bad.line.color.rgb = rgb(ROSE)
        bad.line.width = Pt(0.35)
        disable_shadow(bad)
    reject = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, inch(9.94), inch(3.07), inch(10.13), inch(3.38))
    reject.line.color.rgb = rgb(ROSE)
    reject.line.width = Pt(1.05)
    reject2 = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, inch(10.13), inch(3.07), inch(9.94), inch(3.38))
    reject2.line.color.rgb = rgb(ROSE)
    reject2.line.width = Pt(1.05)
    keep = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, inch(9.36), inch(3.11), inch(0.42), inch(0.27))
    keep.fill.background()
    keep.line.color.rgb = rgb(GREEN)
    keep.line.width = Pt(1.0)
    disable_shadow(keep)
    add_text(slide, "candidate x'", 9.31, 3.90, 0.86, 0.13, size=5.8, color=MUTED)

    gate = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.HEXAGON, inch(10.60), inch(2.98), inch(0.72), inch(0.66))
    gate.fill.solid()
    gate.fill.fore_color.rgb = rgb(WHITE)
    gate.line.color.rgb = rgb(ROSE)
    gate.line.width = Pt(1.15)
    disable_shadow(gate)
    add_text(slide, "Π", 10.76, 3.13, 0.40, 0.18, size=12.5, bold=True, color=ROSE)
    add_text(slide, "project", 10.58, 3.73, 0.76, 0.13, size=5.8, color=ROSE)

    after = add_panel(slide, 11.64, 2.90, 1.05, 0.94, fill=WHITE, line=GREEN, width=0.60, radius_shape=MSO_AUTO_SHAPE_TYPE.RECTANGLE)
    add_image_crop(slide, assets["vine_zoom_02"], 11.72, 2.98, 0.86, 0.58, crop_left=0.08, crop_right=0.08)
    add_text(slide, "committed x*", 11.72, 3.90, 0.86, 0.13, size=5.8, color=GREEN)
    add_arrow(slide, 10.27, 3.31, 10.58, 3.31, color=ROSE, width=0.85)
    add_arrow(slide, 11.34, 3.31, 11.62, 3.31, color=ROSE, width=0.85)
    add_text(slide, "s_{d+1}=(Enc_θ(x*,y), A_{d+1})", 9.54, 4.36, 2.68, 0.18, size=7.0, color=INK, bold=True)

    # Main flow arrows.
    add_arrow(slide, 2.78, 3.30, 3.18, 2.80, color=MUTED, width=1.0)
    add_arrow(slide, 5.54, 2.83, 6.16, 3.24, color=MUTED, width=1.0, dash=True)
    add_arrow(slide, 5.54, 4.12, 6.16, 3.62, color=AMBER, width=1.0)
    add_arrow(slide, 8.60, 3.45, 9.03, 3.45, color=MUTED, width=1.0)

    # Thin operator-use hints, short enough not to masquerade as the data path.
    add_arrow(slide, 7.28, 2.58, 4.92, 1.36, color=GREEN, width=0.55, dash=True)
    add_arrow(slide, 11.20, 2.10, 7.06, 1.36, color=ROSE, width=0.55, dash=True)

    # Commit loop is the main invariant.
    add_elbow_arrow(slide, [(11.90, 4.96), (11.90, 5.74), (1.66, 5.74), (1.66, 4.58)], color=ROSE, width=1.05)
    add_text(slide, "Only the projected, re-encoded state is visible to depth d+1.", 3.64, 5.48, 5.92, 0.24, size=8.6, color=ROSE, bold=True)

    out = OUT_DIR / "ps_rsle_method_overview_20260512.pptx"
    prs.save(out)
    return out


def main() -> None:
    OUT_DIR.mkdir(parents=True, exist_ok=True)
    assets = prepare_component_assets()
    small = build_sparse_interface(assets)
    overview = build_method_overview(assets)
    manifest = {
        "schema": "ps_rsle_method_pptx_20260512",
        "workflow_contract": "Editable PPTX source; PDF/PNG exports are QA and LaTeX inclusion artifacts.",
        "source_docs": [
            "/Users/fanta/Downloads/sparse_latent_generator_interface_figure_brief (1).md",
            "/Users/fanta/Downloads/ps_rsle_ch4_4_1_4_3_revision_notes.md",
        ],
        "outputs": {
            "sparse_latent_generator_interface_pptx": str(small),
            "ps_rsle_method_overview_pptx": str(overview),
        },
        "component_assets": {k: str(v) for k, v in assets.items()},
        "terminology_constraints": {
            "generic_interface_avoids": ["TRELLIS", "TRELLIS2", "O-Voxel", "SLat", "Generative Modelling"],
            "main_method_uses": ["PS-RSLE", "active handles", "generator-control package C_d", "projection + commit"],
        },
    }
    (OUT_DIR / "method_diagram_pptx_manifest_20260512.json").write_text(json.dumps(manifest, indent=2), encoding="utf-8")
    print(json.dumps(manifest, indent=2))


if __name__ == "__main__":
    main()
